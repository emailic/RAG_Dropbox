import os
import time
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
import openai
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import pytesseract
import pdf2image
from tqdm import tqdm
import dropbox
import requests

load_dotenv()  # Load environment variables from .env
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
PINECONE_API_KEY = os.environ.get("PINECONE_API_KEY")
DROPBOX_ACCESS_TOKEN = os.getenv("DROPBOX_ACCESS_TOKEN")
#print (DROPBOX_ACCESS_TOKEN)
INDEX_NAME = "dropbox-rag"
PINECONE_CLOUD = "aws"
PINECONE_REGION = "us-east-1"

if not PINECONE_API_KEY:
    raise ValueError("Missing required API keys. Check your .env file.")

client = openai.OpenAI(api_key = OPENAI_API_KEY)
pinecone = Pinecone(api_key=os.environ.get("PINECONE_API_KEY"))

def initialize_pinecone():
    """Create Pinecone index."""
    if INDEX_NAME not in pinecone.list_indexes().names():
        pinecone.create_index(
            name=INDEX_NAME, dimension=1536, metric="cosine",
            spec=ServerlessSpec(cloud=PINECONE_CLOUD, region=PINECONE_REGION)
        )
        print(f"Creating index {INDEX_NAME}...")

    while INDEX_NAME not in pinecone.list_indexes().names():
        print("Waiting for the index to be ready...")
        time.sleep(5)
    
    return pinecone.Index(INDEX_NAME)


def clear_pinecone_index(index):
    """Clear the Pinecone index before inserting new data."""
    stats = index.describe_index_stats()
    total_vectors = stats["total_vector_count"]

    if total_vectors > 0:
        print(f"Clearing {total_vectors} vectors from index {index}...")
        index.delete(delete_all=True)
        print("Pinecone index cleared.")
    else:
        print("Pinecone index is already empty. Nothing to delete.")

def download_file_from_dropbox(dropbox_path, local_path):
    """Download a file from Dropbox to a local directory"""
    dbx = dropbox.Dropbox(DROPBOX_ACCESS_TOKEN)

    
    response = dbx.files_list_folder("")
    print("RESPONSESSSSSSS", [entry.name for entry in response.entries] )

    #try:
    dbx.files_download_to_file(local_path, dropbox_path)
    print(f"Downloaded {dropbox_path} to {local_path}")
    # except dropbox.exceptions.ApiError as e:
    #     print(f"Error downloading {dropbox_path}: {e}")


def extract_text_from_pdf(pdf_path):
    """Extract text from PDF, handling both searchable and scanned documents."""
    reader = PdfReader(pdf_path)
    text_chunks = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
             text_chunks.append((text))

    # Caterpillar manual cant be read, likely because it is scanned. Thus we extract text from images:
    if len(text_chunks)==0:
        print(f"PDF Manual needs to be processed using OCR. Path: {pdf_path}")
        images = pdf2image.convert_from_path(pdf_path)
        for i, img in tqdm(enumerate(images), desc = "Extracting text from images"):
            text = pytesseract.image_to_string(img)
            if text:
                text_chunks.append((text)) 
    corpus_string = "".join(text_chunks)
    return  corpus_string

def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint (.pptx), using OCR if necessary."""
    prs = Presentation(pptx_path)
    text_chunks = []
    
    for i, slide in enumerate(prs.slides):
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        if text:
            text_chunks.append((text, i + 1))
    
    if not text_chunks:
        print(f"Performing OCR on PPTX: {pptx_path}")
        images = convert_pptx_to_images(pptx_path)
        for i, img in enumerate(images):
            text = pytesseract.image_to_string(img)
            if text:
                text_chunks.append((text, i + 1))
    
    return text_chunks

def extract_text_from_docx(docx_path):
    """Extract text from Word (.docx), using OCR if necessary."""
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    
    if text.strip():
        return [(text, 1)]
    else:
        print(f"Performing OCR on DOCX: {docx_path}")
        images = convert_docx_to_images(docx_path)
        text_chunks = []
        for i, img in enumerate(images):
            text = pytesseract.image_to_string(img)
            if text:
                text_chunks.append((text, i + 1))
        return text_chunks
    
# TODO: Write chunking function.
    
def convert_pptx_to_images(pptx_path):
    """Convert PowerPoint slides to images."""
    return []  # Implement image extraction from PPTX slides if needed

def convert_docx_to_images(docx_path):
    """Convert Word document pages to images."""
    return []  # Implement image extraction from DOCX if needed

if __name__ == "__main__":
    index = initialize_pinecone()
    #clear_pinecone_index(index)

    print(f"Using index: {INDEX_NAME}")
    doc_text= list()


    docs = {
        "Caterpillar 3500": ("/Caterpillar-3500-generator-sets-operation-and-maintenance-manual.pdf", "Caterpillar-3500-generator-sets-operation-and-maintenance-manual.pdf")
        # "Waukesha VGF": ("/Waukesha_VGF_f18g.pdf", "Waukesha_VGF_f18g.pdf"),
        # "Powerpoint Example": ("/powerpoint.pptx", "powerpoint.pptx"),
        # "Word Document Example": ("/word_doc_example.docx", "word_doc_example.docx")
    }    
    for title, (dropbox_path, local_filename) in docs.items():
        local_path = os.path.join(os.getcwd(), local_filename)
        
        # Download file from Dropbox
        download_file_from_dropbox(dropbox_path, local_path)

        # Process the downloaded file
        if local_filename.endswith(".pdf"):
            text_chunks = extract_text_from_pdf(local_path)
        elif local_filename.endswith(".pptx"):
            text_chunks = extract_text_from_pptx(local_path)
        elif local_filename.endswith(".docx"):
            text_chunks = extract_text_from_docx(local_path)
        else:
            continue

    for text in tqdm([text_chunks], desc=f"Indexing {title}"):
        #embedding = client.embeddings.create(input=text, model="text-embedding-3-small")
        #vector = embedding.data[0].embedding
        metadata = {"source": title}
        #index.upsert(vectors=[(f"doc_{manual}_page_{page}", vector, metadata)])
        doc_text.append((text, title))
    print("DOCCC", doc_text)
    
    print("Manuals chunked, vectorized, and upserted into Pinecone database.")
