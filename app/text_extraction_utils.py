import os
import shutil
import tempfile
from typing import List, Dict

from PyPDF2 import PdfReader
import pdf2image
import pytesseract
from docx import Document
import pptx
from PIL import Image
from zipfile import ZipFile
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF, using OCR if needed"""
    texts = []
    
    # try regular text extraction
    logger.info("Trying PdfReader to extract text...")
    reader = PdfReader(pdf_path)
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            texts.append(text.strip())
    
    # else perform OCR
    if not texts:

        images = pdf2image.convert_from_path(pdf_path)
        logger.info("No text retrieved the traditional way. Starting OCR on the PDF file...")
        for img in images:
            text = pytesseract.image_to_string(img).strip()
            if text:
                texts.append(text)
    corpus = "\n".join(texts)
    return corpus

def extract_images_from_docx(docx_path: str, temp_dir: str) -> List[str]:
    """
    Extract images from a DOCX file and save the paths to a temporary directory.
    """
    images = []
    try:
        with ZipFile(docx_path, 'r') as docx_zip:
            for file_name in docx_zip.namelist():
                if file_name.startswith('word/media/'):
                    image_filename = file_name.split('/')[-1]
                    extracted_image_path = os.path.join(temp_dir, image_filename)
                    with docx_zip.open(file_name) as source, open(extracted_image_path, 'wb') as target:
                        target.write(source.read())
                    images.append(extracted_image_path)
    except Exception as e:
        print(f"Error extracting images from DOCX: {e}")
    
    return images



def extract_text_from_docx(docx_path: str) -> str:
    """
    Extract text from DOCX, using OCR if needed.
    """
    texts = []
    images = []
    temp_dir = None
    
    try:
        # try regular text extraction
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                texts.append(text)
        
        # else try OCR
        if not texts:
            temp_dir = tempfile.mkdtemp()
            images = extract_images_from_docx(docx_path, temp_dir)
            for image_path in images:
                try:
                    with Image.open(image_path) as img:
                        text = pytesseract.image_to_string(img).strip()
                        if text:
                            texts.append(text)
                except Exception as img_err:
                    print(f"Error processing image {image_path}: {img_err}")
    
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    
    finally:
        # delete temporary files and directory
        if images and temp_dir:
            for image_path in images:
                if os.path.exists(image_path):
                    os.remove(image_path)
            try:
                shutil.rmtree(temp_dir)
            except Exception as cleanup_err:
                print(f"Error during cleanup: {cleanup_err}")
    
    corpus = "\n".join(texts)
    return corpus

def extract_images_from_pptx(pptx_path: str, temp_dir: str) -> List[str]:
    """
    Extract images from a PowerPoint file and save them to a temporary directory.
    """
    images = []
    try:
        prs = pptx.Presentation(pptx_path)
        
        for i, slide in enumerate(prs.slides):
            # slide.shapes is a sequence of shapes appearing on a slide,
            # first shape in the sequence is the backmost in z-order and the last shape is topmost
            for shape in slide.shapes:
                # check if the shape is a picture
                # MSO_SHAPE_TYPE for pictures. method .shape_type is not impelmented, as I checked,
                # which might be the root of the problem.
                if shape.shape_type == 13:  
                    image = shape.image
                    image_filename = f"slide_{i}_image_{len(images)}.{image.ext}"
                    image_path = os.path.join(temp_dir, image_filename)
                    
                    # Save the image
                    with open(image_path, 'wb') as f:
                        f.write(image.blob)
                    
                    images.append(image_path)
    except Exception as e:
        logger.error(f"Error extracting images from PowerPoint: {e}")
    
    return images


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from PowerPoint file, using OCR for image-based slides.
    """
    texts = []
    images = []
    temp_dir = None
    
    try:
        prs = pptx.Presentation(pptx_path)
        logger.info("Attempting regular text extraction from PowerPoint...")
        for slide in prs.slides:
            # extract text from text boxes
            for shape in slide.shapes:
                # if shape has attribute "text"
                if hasattr(shape, "text"): # we could have also used has_text_frame by pptx
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
                
                # extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                texts.append(cell.text.strip())
        
        # if no text found, try OCR on images
        if not texts:
            temp_dir = tempfile.mkdtemp()
            images = extract_images_from_pptx(pptx_path, temp_dir)
            
            for image_path in images:
                try:
                    with Image.open(image_path) as img:
                        text = pytesseract.image_to_string(img).strip()
                        if text:
                            texts.append(text)
                except Exception as img_err:
                    print(f"Error processing image {image_path}: {img_err}")
    
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint: {e}")
    
    finally:
        # delete temporary files and directory
        if images and temp_dir:
            for image_path in images:
                if os.path.exists(image_path):
                    os.remove(image_path)
            try:
                shutil.rmtree(temp_dir)
            except Exception as cleanup_err:
                logger.error(f"Error during cleanup: {cleanup_err}")
    
    corpus = "\n".join(texts)
    return corpus

def chunk_text(corpus: str, chunk_size: int = 1000) -> List[Dict]:
    """
    Chunk text into smaller pieces with metadata before storing in pinecone.
    """
    chunks = []
    # paragraph based chunking
    paragraphs = [p for p in corpus.split('\n') if p.strip()]
    current_chunk = ""
    
    for para in paragraphs:
        if len(current_chunk) + len(para) < chunk_size:
            current_chunk += para + "\n\n"
        else:
            chunks.append({
                "text": current_chunk.strip(),
                "chunk_id": f"chunk_{len(chunks)+1}"
            })
            current_chunk = para + "\n\n"
    
    if current_chunk.strip():
        chunks.append({
            "text": current_chunk.strip(),
            "chunk_id": f"chunk_{len(chunks)+1}"
        })

    return chunks
    
#print(chunk_text(extract_text_from_pptx("example.pptx")))
